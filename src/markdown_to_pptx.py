from pptx import Presentation
from pptx.util import Inches, Pt
import os


def markdown_to_pptx(markdown_text: str) -> str:
    """将 Markdown 文本转换为 PowerPoint 文件"""
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    lines = markdown_text.split("\n")

    for line in lines:
        if line.strip() == "":
            continue
        slide = prs.slides.add_slide(slide_layout)
        title, content = line.split(" ", 1) if " " in line else ("", line)

        # 设置标题
        title_placeholder = slide.shapes.title
        title_placeholder.text = title

        # 设置正文
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = content

    # 保存 PowerPoint 文件
    output_path = "/outputs/ChatPPT_Output.pptx"
    prs.save(output_path)
    return output_path
